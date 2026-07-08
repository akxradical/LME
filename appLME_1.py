"""
Commodities Trading Terminal
Cu · Al · Brent · Silver | Market Intel | PPTX Export
"""
import io, xml.etree.ElementTree as ET
from datetime import datetime
from urllib.request import urlopen, Request

import gspread, numpy as np, pandas as pd
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt, matplotlib.ticker as mticker
from matplotlib.lines import Line2D
import plotly.graph_objects as go, plotly.express as px
import streamlit as st
from google.oauth2.service_account import Credentials
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SPREADSHEET_ID = "1zLyMANY56oFRwFug04WYavGUH_NAlRH8M3c-TXIRDlI"
SHEETS = {"cu":"LME Copper","al":"LME Aluminium","br":"Brent Oil","ag":"LBMA Silver"}
SCOPES = ["https://www.googleapis.com/auth/spreadsheets.readonly",
          "https://www.googleapis.com/auth/drive.readonly"]
BG="#0D1117";CARD="#161B22";CARD2="#1C2128";BD="#30363D";GD="#21262D"
GR="#00C853";RD="#FF1744";GD2="#FFD600";BL="#2979FF"
TX="#E6EDF3";SC="#8B949E";MT="#484F58"
CLR={"cu":"#E87B35","al":"#78909C","br":"#4CAF50","ag":"#C0C0C0"}
NAMES={"cu":"Copper (Cu)","al":"Aluminium (Al)","br":"Brent Crude Oil","ag":"Silver (Ag)"}
UNITS={"cu":"USD/MT","al":"USD/MT","br":"USD/BBL","ag":"USD/OZ"}

st.set_page_config(page_title="Commodities Terminal",page_icon="📊",layout="wide",initial_sidebar_state="expanded")
st.markdown("""<style>
html,body,[data-testid="stAppViewContainer"],[data-testid="stApp"]{background:#0D1117!important;color:#E6EDF3;font-family:'JetBrains Mono',monospace}
[data-testid="stSidebar"]{background:#161B22!important;border-right:1px solid #30363D}
[data-testid="stHeader"]{background:transparent!important}
.block-container{padding-top:1rem;max-width:100%}
.tc{background:#161B22;border:1px solid #30363D;border-radius:8px;padding:12px 16px;margin-bottom:8px}
.tl{font-size:9px;letter-spacing:1.2px;color:#8B949E;text-transform:uppercase;margin-bottom:3px}
.tv{font-size:24px;font-weight:700;color:#E6EDF3}.ts{font-size:10px;margin-top:2px}
.up{color:#00C853}.dn{color:#FF1744}
.sh{font-size:11px;letter-spacing:1.5px;color:#8B949E;text-transform:uppercase;border-bottom:1px solid #30363D;padding-bottom:6px;margin-bottom:12px}
.dt{width:100%;border-collapse:collapse;font-size:12px}
.dt th{background:#1C2128;color:#8B949E;font-size:10px;letter-spacing:1px;text-transform:uppercase;padding:8px 12px;text-align:right;border-bottom:1px solid #30363D}
.dt th:first-child{text-align:left}.dt td{padding:7px 12px;border-bottom:1px solid #1C2128;text-align:right;color:#E6EDF3}
.dt td:first-child{text-align:left;color:#8B949E}.dt tr:hover td{background:#1C2128}
.stTabs [data-baseweb="tab-list"]{background:#161B22;border-radius:8px;padding:4px;gap:3px;border:1px solid #30363D}
.stTabs [data-baseweb="tab"]{background:transparent;color:#8B949E;border-radius:6px;padding:7px 14px;font-size:11px}
.stTabs [aria-selected="true"]{background:#1C2128!important;color:#E6EDF3!important;border:1px solid #30363D!important}
.stButton>button{background:#2979FF;color:white;border:none;border-radius:6px;font-weight:600;width:100%;padding:10px}
[data-testid="stDownloadButton"]>button{background:#1B5E20;color:#00C853;border:1px solid #00C853;border-radius:6px;font-weight:600;width:100%;padding:10px}
.pc{background:#161B22;border:1px solid #30363D;border-radius:10px;padding:14px;margin-bottom:10px}
.ni{background:#161B22;border-left:3px solid #2979FF;padding:10px 14px;margin-bottom:8px;border-radius:0 6px 6px 0}
.nt{color:#E6EDF3;font-size:13px;font-weight:600;text-decoration:none}.nt:hover{color:#2979FF}
.ns{font-size:10px;color:#484F58;margin-top:3px}
.ev{background:#161B22;border:1px solid #30363D;border-radius:8px;padding:10px 14px;margin-bottom:6px;display:flex;gap:12px;align-items:center}
.ev-date{font-size:11px;font-weight:700;color:#FFD600;min-width:60px}
.ev-body{font-size:11px;color:#8B949E}
.ev-tag{font-size:9px;padding:2px 6px;border-radius:3px;font-weight:600}
#MainMenu,footer,header{visibility:hidden}[data-testid="stToolbar"]{display:none}
</style>""",unsafe_allow_html=True)

# ═══ DATA ═══
@st.cache_data(ttl=300)
def load(name):
    try:
        info=dict(st.secrets["gcp_service_account"]); creds=Credentials.from_service_account_info(info,scopes=SCOPES)
        gc=gspread.authorize(creds); ws=gc.open_by_key(st.secrets.get("SPREADSHEET_ID",SPREADSHEET_ID)).worksheet(name)
        df=pd.DataFrame(ws.get_all_records()); df.columns=[c.strip() for c in df.columns]
        dc=next(c for c in df.columns if "date" in c.lower()); pc=next(c for c in df.columns if "price" in c.lower() or "cash" in c.lower())
        df=df.rename(columns={dc:"Date",pc:"Price"})
        df["Date"]=pd.to_datetime(df["Date"],dayfirst=True,errors="coerce")
        df["Price"]=pd.to_numeric(df["Price"].astype(str).str.replace(",",""),errors="coerce")
        return df.dropna(subset=["Date","Price"]).sort_values("Date").reset_index(drop=True)
    except Exception:
        # NEVER return fake data — return empty so the app shows an honest "no data" state
        return pd.DataFrame(columns=["Date","Price"])

@st.cache_data(ttl=3600)
def get_usd_inr():
    try:
        url="https://fred.stlouisfed.org/graph/fredgraph.csv?id=DEXINUS&observation_start=2025-01-01"
        req=Request(url,headers={"User-Agent":"Mozilla/5.0"})
        with urlopen(req,timeout=5) as r: lines=r.read().decode().strip().split("\n")
        for line in reversed(lines[1:]):
            parts=line.strip().split(",")
            if len(parts)>=2 and parts[1].strip() not in (".",""):
                return float(parts[1].strip())
    except: pass
    return 85.5  # fallback

# ═══ ANALYTICS ═══
def mavg(df):
    d=df.copy();d["YM"]=d["Date"].dt.to_period("M");m=d.groupby("YM")["Price"].mean().reset_index()
    m["Date"]=m["YM"].dt.to_timestamp();m["Label"]=m["YM"].dt.strftime("%b-%y");m["Price"]=m["Price"].round(2);return m.drop(columns="YM")
def qavg(df):
    d=df.copy();d["YQ"]=d["Date"].dt.to_period("Q");q=d.groupby("YQ")["Price"].mean().reset_index()
    q["Date"]=q["YQ"].dt.to_timestamp();q["Label"]=q["YQ"].dt.strftime("Q%q-%Y");q["Price"]=q["Price"].round(2);return q.drop(columns="YQ")
def ravg(df,w=3): d=df.copy().sort_values("Date");d[f"MA{w}"]=d["Price"].rolling(w).mean().round(2);return d
def lfc(m,p=3):
    if len(m)<2:return pd.DataFrame(columns=["Date","Price","Label"])
    x=np.arange(len(m));s,i=np.polyfit(x,m["Price"].values,1);last=m["Date"].iloc[-1]
    fd=[last+pd.DateOffset(months=k+1) for k in range(p)];fp=[i+s*(len(m)+k+1) for k in range(p)]
    return pd.DataFrame({"Date":fd,"Price":np.round(fp,2),"Label":[d.strftime("%b-%y") for d in fd]})

# ═══ DATA-DRIVEN INSIGHT ENGINE (no LLM — pure statistics) ═══
def analyze_price(df, unit="USD/MT"):
    """Generate quantitative insights from raw price data — works for any commodity."""
    d = df.copy().sort_values("Date").reset_index(drop=True)
    if len(d) < 30:
        return []
    price = d["Price"]
    cur = price.iloc[-1]

    # Windows
    p7   = price.iloc[-7]  if len(d) > 7  else price.iloc[0]
    p30  = price.iloc[-30] if len(d) > 30 else price.iloc[0]
    p90  = price.iloc[-90] if len(d) > 90 else price.iloc[0]
    ytd  = price.iloc[0]

    chg7  = (cur - p7) / p7 * 100
    chg30 = (cur - p30) / p30 * 100
    chg90 = (cur - p90) / p90 * 100
    chg_ytd = (cur - ytd) / ytd * 100

    # Moving averages
    ma20 = price.rolling(20).mean().iloc[-1]
    ma50 = price.rolling(50).mean().iloc[-1] if len(d) >= 50 else ma20

    # Volatility (30d annualized-ish)
    vol30 = price.pct_change().tail(30).std() * 100

    # 52-week / full-period high-low
    hi = price.max(); lo = price.min()
    pos_in_range = (cur - lo) / (hi - lo) * 100 if hi > lo else 50

    insights = []

    # 1. TREND direction (based on 30d momentum + MA position)
    if chg30 > 3 and cur > ma20:
        trend_txt = f"Prices are in a clear <b style='color:#00C853'>uptrend</b> — up {abs(chg30):.1f}% over 30 days, trading above the 20-day average (${ma20:,.2f}). Momentum favors buyers."
        tag, tc = "RISING", "#00C853"
    elif chg30 < -3 and cur < ma20:
        trend_txt = f"Prices are in a <b style='color:#FF1744'>downtrend</b> — down {abs(chg30):.1f}% over 30 days, trading below the 20-day average (${ma20:,.2f}). Selling pressure dominates."
        tag, tc = "FALLING", "#FF1744"
    else:
        trend_txt = f"Prices are <b style='color:#FFD600'>consolidating</b> — moving {chg30:+.1f}% over 30 days near the 20-day average (${ma20:,.2f}). Market awaiting a catalyst."
        tag, tc = "STABLE", "#FFD600"
    insights.append(("Trend", trend_txt, tag, tc))

    # 2. MOMENTUM (short vs medium term)
    if chg7 > 0 and chg30 > 0:
        mom = f"Short-term momentum is <b style='color:#00C853'>positive</b>: +{chg7:.1f}% this week reinforces the +{chg30:.1f}% monthly gain. Rally has staying power."
    elif chg7 < 0 and chg30 > 0:
        mom = f"Momentum is <b style='color:#FFD600'>cooling</b>: down {abs(chg7):.1f}% this week despite +{chg30:.1f}% for the month — a possible pullback or profit-taking."
    elif chg7 > 0 and chg30 < 0:
        mom = f"Early <b style='color:#00C853'>reversal signal</b>: +{chg7:.1f}% this week against a -{abs(chg30):.1f}% monthly decline — downtrend may be bottoming."
    else:
        mom = f"Momentum is <b style='color:#FF1744'>weak</b>: down {abs(chg7):.1f}% this week and {abs(chg30):.1f}% this month — sustained selling."
    insights.append(("Momentum", mom, None, None))

    # 3. VOLATILITY regime
    if vol30 > 2.5:
        vtxt = f"<b style='color:#FF1744'>High volatility</b> ({vol30:.1f}% daily swings). Prices are choppy — consider staggered buying to average out risk rather than a single large order."
    elif vol30 < 1.0:
        vtxt = f"<b style='color:#00C853'>Low volatility</b> ({vol30:.1f}% daily swings). Calm market — favorable window to lock in forward contracts at predictable prices."
    else:
        vtxt = f"<b style='color:#FFD600'>Moderate volatility</b> ({vol30:.1f}% daily swings). Normal market conditions for procurement timing."
    insights.append(("Volatility", vtxt, None, None))

    # 4. POSITION in range (support/resistance)
    if pos_in_range > 80:
        rtxt = f"Trading near <b style='color:#FF1744'>period highs</b> ({pos_in_range:.0f}% of range, ${lo:,.0f}–${hi:,.0f}). Limited upside room; risk of pullback. Not an ideal entry for large buys."
    elif pos_in_range < 20:
        rtxt = f"Trading near <b style='color:#00C853'>period lows</b> ({pos_in_range:.0f}% of range, ${lo:,.0f}–${hi:,.0f}). Attractive entry zone if fundamentals support a rebound."
    else:
        rtxt = f"Trading in the <b style='color:#FFD600'>mid-range</b> ({pos_in_range:.0f}% of ${lo:,.0f}–${hi:,.0f}). Neither cheap nor expensive versus recent history."
    insights.append(("Price Position", rtxt, None, None))

    # 5. MA crossover (golden/death cross signal)
    if len(d) >= 50:
        if ma20 > ma50 and chg30 > 0:
            cxt = f"The 20-day average (${ma20:,.2f}) sits above the 50-day (${ma50:,.2f}) — a <b style='color:#00C853'>bullish structure</b> confirming the uptrend."
        elif ma20 < ma50:
            cxt = f"The 20-day average (${ma20:,.2f}) sits below the 50-day (${ma50:,.2f}) — a <b style='color:#FF1744'>bearish structure</b>; caution on the medium-term outlook."
        else:
            cxt = f"The 20-day (${ma20:,.2f}) and 50-day (${ma50:,.2f}) averages are converging — a <b style='color:#FFD600'>trend change</b> may be forming."
        insights.append(("Trend Structure", cxt, None, None))

    return insights, {"chg7":chg7,"chg30":chg30,"chg90":chg90,"ytd":chg_ytd,"vol":vol30,"pos":pos_in_range,"tag":tag,"tc":tc}


def render_data_insights(df, unit="USD/MT"):
    """Render the computed insights as CEO-ready cards."""
    result = analyze_price(df, unit)
    if not result:
        st.caption("Not enough data for analysis (need 30+ days).")
        return
    insights, stats = result

    # Summary banner
    st.markdown(
        '<div style="background:linear-gradient(90deg,'+stats["tc"]+'22,transparent);'
        'border-left:4px solid '+stats["tc"]+';border-radius:6px;padding:12px 16px;margin-bottom:14px">'
        '<span style="font-size:11px;font-weight:700;color:'+stats["tc"]+';letter-spacing:1px">● '+stats["tag"]+'</span>'
        '<div style="font-size:11px;color:#8B949E;margin-top:6px;line-height:1.5">'
        '7-Day: <b style="color:'+("#00C853" if stats["chg7"]>=0 else "#FF1744")+'">'+f'{stats["chg7"]:+.1f}%</b> &nbsp;·&nbsp; '
        '30-Day: <b style="color:'+("#00C853" if stats["chg30"]>=0 else "#FF1744")+'">'+f'{stats["chg30"]:+.1f}%</b> &nbsp;·&nbsp; '
        '90-Day: <b style="color:'+("#00C853" if stats["chg90"]>=0 else "#FF1744")+'">'+f'{stats["chg90"]:+.1f}%</b> &nbsp;·&nbsp; '
        'YTD: <b style="color:'+("#00C853" if stats["ytd"]>=0 else "#FF1744")+'">'+f'{stats["ytd"]:+.1f}%</b>'
        '</div></div>',
        unsafe_allow_html=True)

    # Individual insight cards
    for item in insights:
        title, txt = item[0], item[1]
        st.markdown(
            '<div class="pc" style="margin-bottom:8px">'
            '<div style="font-size:11px;font-weight:700;color:#E6EDF3;margin-bottom:5px">▸ '+title+'</div>'
            '<div style="font-size:11px;color:#8B949E;line-height:1.6">'+txt+'</div></div>',
            unsafe_allow_html=True)

# ═══ CHARTS ═══
_CL=dict(paper_bgcolor=CARD,plot_bgcolor=CARD,font=dict(color=TX,family="monospace",size=11),
    margin=dict(l=55,r=15,t=65,b=55),xaxis=dict(showgrid=False,zeroline=False,color=SC,linecolor=BD),
    yaxis=dict(gridcolor=GD,zeroline=False,color=SC,linecolor=BD,tickformat=","),
    legend=dict(bgcolor="rgba(0,0,0,0)",bordercolor=BD,orientation="h",y=-0.18,font=dict(size=10)))

def chart_live(raw,metal,color,unit="USD/MT"):
    df=raw.copy().sort_values("Date").reset_index(drop=True);r3=ravg(df,3)
    cur=df["Price"].iloc[-1];prev=df["Price"].iloc[-2] if len(df)>1 else cur
    ch=cur-prev;pc=ch/prev*100 if prev else 0;ar="▲" if ch>=0 else "▼";cl=GR if ch>=0 else RD
    r,g,b=bytes.fromhex(color.lstrip("#"))
    fig=go.Figure()
    fig.add_trace(go.Scatter(x=df["Date"],y=df["Price"],fill="tozeroy",fillcolor=f"rgba({r},{g},{b},0.07)",
        line=dict(color=color,width=1.5),mode="lines",name="Daily",hovertemplate="<b>%{x|%d %b %Y}</b><br>$%{y:,.2f}<extra></extra>"))
    fig.add_trace(go.Scatter(x=r3["Date"],y=r3["MA3"],line=dict(color=GD2,width=1.2,dash="dot"),mode="lines",name="3M MA"))
    fl=df["Date"].iloc[0].strftime("%d %b %Y");ll=df["Date"].iloc[-1].strftime("%d %b %Y")
    _b={k:v for k,v in _CL.items() if k not in ("xaxis","yaxis")}
    fig.update_layout(**_b,title=dict(text=f"<b>{metal}</b>  <span style='font-size:22px;font-weight:700;color:{cl}'>${cur:,.2f}</span>  "
        f"<span style='font-size:13px;color:{cl}'>{ar} ${abs(ch):,.2f} ({ar}{abs(pc):.2f}%)</span>"
        f"<br><span style='font-size:10px;color:{SC}'>{unit} · {fl} → {ll}</span>",font=dict(size=15,color=TX)),
        xaxis=dict(showgrid=False,zeroline=False,color=SC,linecolor=BD,
            rangeslider=dict(visible=True,bgcolor=BG,thickness=0.04),
            rangeselector=dict(bgcolor=CARD2,activecolor=BD,font=dict(color=SC,size=10),
                buttons=[dict(count=1,label="1M",step="month",stepmode="backward"),
                    dict(count=3,label="3M",step="month",stepmode="backward"),
                    dict(count=6,label="6M",step="month",stepmode="backward"),
                    dict(count=1,label="YTD",step="year",stepmode="todate"),dict(step="all",label="ALL")])),
        yaxis=dict(gridcolor=GD,zeroline=False,color=SC,linecolor=BD,tickformat=",",side="right"),height=420,hovermode="x unified")
    return fig

def chart_q(q,fq,metal,color):
    q2=q.copy();q2["p"]=q2["Price"].shift(1);q2["u"]=q2["Price"]>=q2["p"]
    fig=go.Figure()
    fig.add_trace(go.Bar(x=q2["Label"],y=q2["Price"],marker_color=[GR if u else RD for u in q2["u"]],opacity=0.85,
        text=[f"{p:,.0f}" for p in q2["Price"]],textposition="outside",textfont=dict(size=9,color=TX),name="Qtr Avg"))
    if len(fq):
        fig.add_trace(go.Bar(x=fq["Label"],y=fq["Price"],marker_color=GD2,opacity=0.5,
            text=[f"{p:,.0f}" for p in fq["Price"]],textposition="outside",textfont=dict(size=9,color=GD2),name="Forecast"))
    fig.update_layout(**_CL,title=dict(text=f"<b>{metal}</b> <span style='color:{SC};font-size:12px'>Quarterly</span>",font=dict(size=15)),height=350,barmode="group")
    return fig

def chart_roll(raw,metal,color):
    r3=ravg(raw,3);r6=ravg(raw,6);fig=go.Figure()
    fig.add_trace(go.Scatter(x=r3["Date"],y=r3["Price"],line=dict(color=BD,width=0.8),mode="lines",name="Daily",opacity=0.5))
    fig.add_trace(go.Scatter(x=r3["Date"],y=r3["MA3"],line=dict(color=color,width=2),mode="lines",name="3M"))
    fig.add_trace(go.Scatter(x=r6["Date"],y=r6["MA6"],line=dict(color=GD2,width=1.5,dash="dash"),mode="lines",name="6M"))
    fig.update_layout(**_CL,title=dict(text=f"<b>{metal}</b> <span style='color:{SC};font-size:12px'>Rolling Avg</span>",font=dict(size=15)),height=350)
    return fig

# ═══ KPIS ═══
def kpis(m,unit="USD/MT"):
    c=m["Price"].iloc[-1];cl=m["Label"].iloc[-1];lo=m["Price"].min();ll=m.loc[m["Price"].idxmin(),"Label"]
    hi=m["Price"].max();hl=m.loc[m["Price"].idxmax(),"Label"]
    p=m["Price"].iloc[-2] if len(m)>1 else c;ch=c-p;pc=ch/p*100 if p else 0
    f=m["Price"].iloc[0];yc=c-f;yp=yc/f*100 if f else 0
    ud="up" if ch>=0 else "dn";ya="up" if yc>=0 else "dn"
    ar="▲" if ch>=0 else "▼";yr="▲" if yc>=0 else "▼"
    st.markdown(
        '<div class="tc"><div class="tl">CURRENT · '+cl+'</div><div class="tv">$'+f'{c:,.2f}'+'</div>'
        '<div class="ts '+ud+'">'+ar+' $'+f'{abs(ch):,.2f}'+' ('+f'{abs(pc):.1f}'+'%) MoM</div></div>'
        '<div class="tc"><div class="tl">HIGH · '+hl+'</div><div class="tv up">$'+f'{hi:,.2f}'+'</div></div>'
        '<div class="tc"><div class="tl">LOW · '+ll+'</div><div class="tv dn">$'+f'{lo:,.2f}'+'</div></div>'
        '<div class="tc"><div class="tl">YTD</div><div class="tv '+ya+'">'+yr+' '+f'{abs(yp):.1f}'+'%</div>'
        '<div class="ts '+ya+'">'+yr+' $'+f'{abs(yc):,.0f}'+' since '+m["Label"].iloc[0]+'</div></div>',
        unsafe_allow_html=True)

# ═══ TABLE ═══
def tbl(monthly):
    m2=monthly.copy();m2["Ch"]=m2["Price"].diff().round(2);m2["P"]=(m2["Price"].pct_change()*100).round(2);m2["A3"]=m2["Price"].rolling(3).mean().round(2)
    rows=""
    for _,r in m2.iterrows():
        ch="" if pd.isna(r["Ch"]) else '<span class="'+("up" if r["Ch"]>=0 else "dn")+'">'+("+" if r["Ch"]>=0 else "")+f'{r["Ch"]:,.2f}</span>'
        pc="" if pd.isna(r["P"]) else '<span class="'+("up" if r["P"]>=0 else "dn")+'">'+("+" if r["P"]>=0 else "")+f'{r["P"]:.1f}%</span>'
        a3="—" if pd.isna(r["A3"]) else f'${r["A3"]:,.2f}'
        rows+='<tr><td>'+r["Label"]+'</td><td>$'+f'{r["Price"]:,.2f}'+'</td><td>'+ch+'</td><td>'+pc+'</td><td style="color:#8B949E">'+a3+'</td></tr>'
    st.markdown('<table class="dt"><thead><tr><th>Month</th><th>Price</th><th>MoM</th><th>%</th><th>3M Avg</th></tr></thead><tbody>'+rows+'</tbody></table>',unsafe_allow_html=True)

# ═══ TICKER BAR ═══
def ticker_bar(data):
    def blk(name,m,clr):
        c=m["Price"].iloc[-1];p=m["Price"].iloc[-2] if len(m)>1 else c;ch=c-p;pc=ch/p*100 if p else 0
        badge="lme-cup" if ch>=0 else "lme-cdn";arr="&#9650;" if ch>=0 else "&#9660;"
        return ('<div class="lme-pb"><div class="lme-mn">'+name+'</div>'
            '<div class="lme-pr"><span class="lme-pv" style="color:'+clr+'">'+'${:,.2f}'.format(c)+'</span></div>'
            '<div class="lme-cr"><span class="lme-badge '+badge+'">'+arr+' {:.2f}%'.format(abs(pc))+'</span>'
            '<span class="lme-ab">'+arr+' ${:,.2f}'.format(abs(ch))+' MoM</span></div></div>')
    now=datetime.now().strftime("%d %b %Y %H:%M")
    css=('<style>.lme-wrap{background:#0A0F1E;border-bottom:2px solid #B87333;margin-bottom:12px}'
        '.lme-top{background:#050A14;padding:7px 24px;display:flex;align-items:center;justify-content:space-between;border-bottom:1px solid #1C2333}'
        '.lme-logo{font-size:22px;font-weight:900;letter-spacing:3px;color:#FFF;font-family:"Arial Black",Arial,sans-serif}'
        '.lme-bar{background:#0A0F1E;padding:12px 20px;display:flex;align-items:stretch;border-bottom:1px solid #1C2333;flex-wrap:wrap}'
        '.lme-pb{display:flex;flex-direction:column;padding:6px 24px 6px 0;margin-right:20px;border-right:1px solid #1C2333;min-width:170px}'
        '.lme-mn{font-size:9px;letter-spacing:1.5px;color:#8B949E;text-transform:uppercase;margin-bottom:4px}'
        '.lme-pr{display:flex;align-items:baseline}.lme-pv{font-size:22px;font-weight:700;letter-spacing:-0.5px}'
        '.lme-cr{display:flex;align-items:center;gap:6px;margin-top:3px}'
        '.lme-badge{font-size:10px;font-weight:600;padding:2px 7px;border-radius:3px}'
        '.lme-cup{background:rgba(0,200,83,0.15);color:#00C853}.lme-cdn{background:rgba(255,23,68,0.15);color:#FF1744}'
        '.lme-ab{font-size:9px;color:#8B949E}'
        '.lme-rf{margin-left:auto;display:flex;flex-direction:column;justify-content:center;padding-left:20px;border-left:1px solid #1C2333}'
        '.lme-ld{display:inline-block;width:7px;height:7px;background:#00C853;border-radius:50%;margin-right:5px;vertical-align:middle;animation:lp 2s infinite}'
        '@keyframes lp{0%,100%{opacity:1}50%{opacity:0.3}}</style>')
    labels={"cu":"COPPER &middot; USD/MT","al":"ALUMINIUM &middot; USD/MT",
            "br":"BRENT &middot; USD/BBL","ag":"SILVER &middot; USD/OZ"}
    blocks=""
    for key in ["cu","al","br","ag"]:
        if key+"_m" in data:
            blocks+=blk(labels[key],data[key+"_m"],CLR[key])

    body=('<div class="lme-wrap"><div class="lme-top"><div>'
        '<div class="lme-logo">LME<span style="color:#B87333">.</span></div>'
        '<div style="font-size:9px;color:#8B949E;letter-spacing:2px;text-transform:uppercase;margin-top:1px">Commodities Terminal</div>'
        '</div><div style="font-size:10px;color:#484F58;letter-spacing:1px">DAILY SETTLEMENT</div></div><div class="lme-bar">'
        +blocks
        +'<div class="lme-rf"><span style="font-size:9px;color:#484F58;letter-spacing:1.5px;text-transform:uppercase">'
        '<span class="lme-ld"></span>Live</span>'
        '<span style="font-size:12px;color:#8B949E;margin-top:2px">'+now+'</span>'
        '<span style="font-size:9px;color:#484F58;margin-top:1px">Auto 5 min</span></div>'
        '</div></div><br>')
    st.markdown(css+body,unsafe_allow_html=True)

# ═══ MARKET INTEL DATA ═══
TOP5={"Copper":[("🇨🇱","Chile","5.8M MT","Escondida, Collahuasi — world #1, ~27% global supply"),
    ("🇵🇪","Peru","2.7M MT","Cerro Verde, Antamina — vast reserves, needs H₂SO₄ from China for SX-EW leaching"),
    ("🇨🇩","DR Congo","2.5M MT","Kamoa-Kakula — fastest growing; political instability risk"),
    ("🇨🇳","China","1.9M MT","Largest smelter & consumer; imports 70%+ concentrate; controls H₂SO₄ chain"),
    ("🇺🇸","USA","1.1M MT","Morenci (Freeport) — declining grades, permit delays")],
"Aluminium":[("🇨🇳","China","41M MT","60% of world output; Yunnan hydropower curtailments in dry season"),
    ("🇮🇳","India","4.1M MT","Hindalco, Vedanta — cheap coal power; growing exports"),
    ("🇷🇺","Russia","3.8M MT","Rusal — sanctions risk; LME warehouse ban discussions"),
    ("🇨🇦","Canada","3.1M MT","Rio Tinto — 100% hydropower 'green aluminium' premium"),
    ("🇦🇪","UAE","2.7M MT","EGA — gas-powered Gulf hub for Asian/African bauxite")],
"Brent Oil":[("🇺🇸","USA","13.3M bpd","Largest producer; Permian shale; SPR releases influence"),
    ("🇸🇦","Saudi Arabia","12.5M bpd","OPEC+ swing producer; Aramco spare capacity"),
    ("🇷🇺","Russia","10.8M bpd","Urals discount; EU embargo rerouted to India/China"),
    ("🇮🇶","Iraq","4.5M bpd","Basra crude tied to Brent; OPEC quota compliance varies"),
    ("🇮🇷","Iran","3.4M bpd","Sanctions limit exports; shadow fleet; nuclear deal talks")],
"Silver":[("🇲🇽","Mexico","6,300 MT","World #1 — Fresnillo, First Majestic; mining reform risk under AMLO successor"),
    ("🇨🇳","China","3,400 MT","#2 producer AND largest consumer (solar PV, electronics)"),
    ("🇵🇪","Peru","3,100 MT","Antamina, Buenaventura — community protests disrupt output"),
    ("🇨🇱","Chile","1,600 MT","Byproduct of copper mining — output tied to Cu capex cycles"),
    ("🇵🇱","Poland","1,300 MT","KGHM — Europe's sole major silver-copper mine; energy cost pressure")]}

INSIGHTS={"Copper":[("H₂SO₄ Supply Chain","Peru's SX-EW leaching requires sulfuric acid from Chinese smelters. Environmental shutdowns in China directly constrain Peru's 2.7M MT output."),
    ("EV Copper Intensity","Each EV needs ~83kg copper vs 23kg for ICE. IEA projects 6-8M MT supply deficit by 2030 without new mines (10-15 year lead time)."),
    ("TC/RC Squeeze","Smelter treatment charges at multi-year lows — excess capacity chasing scarce concentrate.")],
"Aluminium":[("Yunnan Dry Season","12% of global smelting in Yunnan (hydro-dependent). Nov-Apr curtailments remove 1-2M MT, causing seasonal spikes."),
    ("EU CBAM 2026","Carbon Border Tax adds €50-100/MT to carbon-intensive imports. Benefits Canadian/Nordic green smelters."),
    ("Russia Overhang","Rusal = 6% global supply. LME ban discussions create volatility; sanctions reshape trade flows.")],
"Brent Oil":[("OPEC+ Cuts","Saudi voluntary cuts ~2.2M bpd support $70-80 floor. Iraq/Kazakhstan chronic overproducers."),
    ("Red Sea Disruption","Houthi attacks reroute 12% of global trade via Cape — adds $1-2M per voyage, $5-10/bbl risk premium."),
    ("Shale Ceiling","Permian breakeven ~$45-55/bbl. Above $70, US rigs increase in 3-6 months — automatic price ceiling.")],
"Silver":[("Solar Demand Surge","Each GW of solar PV uses ~20 tonnes of silver. Global installs hit 400GW+ in 2025. Silver demand from solar doubled since 2020."),
    ("India Import Duty","India cut silver import duty from 15% to 6% in July 2024 — imports surged 40%. Any reversal would crash demand."),
    ("Mexico Mining Reform","AMLO's successor continues mining concession freeze. New permits halted — constrains world's #1 producer."),
    ("Gold-Silver Ratio","Currently ~85:1 vs historical 60:1. Silver considered undervalued relative to gold — attracts momentum traders.")]}

EVENTS=[
    ("Jun 2026","OPEC+ Meeting","Brent","rgba(76,175,80,0.15)","#4CAF50","Production quota decision for H2 2026"),
    ("Jun 2026","US FOMC / Fed","All","rgba(41,121,255,0.15)","#2979FF","Rate decision — impacts USD, all commodity prices"),
    ("Jun 2026","China PMI Release","Cu/Al","rgba(232,123,53,0.15)","#E87B35","Manufacturing activity gauge — demand signal"),
    ("Jul 2026","LME Warehouse Report","Cu/Al","rgba(232,123,53,0.15)","#E87B35","Inventory levels signal supply tightness"),
    ("Jul 2026","India Union Budget","Ag","rgba(192,192,192,0.15)","#C0C0C0","Import duty changes on silver/gold"),
    ("Aug 2026","Jackson Hole Symposium","All","rgba(41,121,255,0.15)","#2979FF","Fed signals on rates — major USD mover"),
    ("Sep 2026","OPEC+ Meeting","Brent","rgba(76,175,80,0.15)","#4CAF50","Q4 production targets"),
    ("Oct 2026","China Party Congress","Cu/Al/Ag","rgba(232,123,53,0.15)","#E87B35","Infrastructure stimulus signals"),
]

def render_top5(c):
    items=TOP5.get(c,[]);html='<div style="display:grid;gap:8px">'
    for i,(flag,country,vol,note) in enumerate(items):
        rc=["#FFD600","#C0C0C0","#CD7F32","#8B949E","#484F58"][i]
        html+=('<div class="pc" style="display:flex;gap:12px;align-items:flex-start">'
            '<div style="font-size:12px;font-weight:800;color:'+rc+';min-width:22px">#'+str(i+1)+'</div>'
            '<div style="font-size:26px;line-height:1">'+flag+'</div><div style="flex:1">'
            '<div style="font-size:13px;font-weight:700;color:#E6EDF3">'+country
            +'<span style="font-size:10px;color:#8B949E;margin-left:8px">'+vol+'</span></div>'
            '<div style="font-size:11px;color:#8B949E;margin-top:3px;line-height:1.5">'+note+'</div></div></div>')
    st.markdown(html+'</div>',unsafe_allow_html=True)

def render_insights(c):
    for t,txt in INSIGHTS.get(c,[]):
        st.markdown('<div class="pc"><div style="font-size:11px;font-weight:700;color:#FFD600;margin-bottom:5px">⚡ '+t+'</div>'
            '<div style="font-size:11px;color:#8B949E;line-height:1.6">'+txt+'</div></div>',unsafe_allow_html=True)

@st.cache_data(ttl=1800)
def news(q,n=6):
    url="https://news.google.com/rss/search?q="+q.replace(" ","+")+"+price&hl=en&gl=US&ceid=US:en"
    try:
        req=Request(url,headers={"User-Agent":"Mozilla/5.0"})
        with urlopen(req,timeout=5) as r: tree=ET.parse(r)
        items=[]
        for i in tree.findall(".//item")[:n]:
            t=i.findtext("title","") or ""
            l=i.findtext("link","") or ""
            s=i.findtext("source","") or ""
            d=i.findtext("pubDate","") or ""
            t=t.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace('"',"&quot;")
            s=s.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            if t: items.append({"t":t,"l":l,"s":s,"d":d})
        return items
    except: return []

def render_news(c):
    qm={"Copper":"copper commodity LME","Aluminium":"aluminium commodity LME","Brent Oil":"brent crude oil OPEC","Silver":"silver LBMA precious metal"}
    try: items=news(qm.get(c,c))
    except: items=[]
    if not items:
        st.caption("No news available right now.")
        return
    for i in items:
        try:
            st.markdown('<div class="ni"><a class="nt" href="'+str(i.get("l",""))+'" target="_blank">'+str(i.get("t",""))+'</a>'
                '<div class="ns">'+str(i.get("s",""))+' · '+str(i.get("d",""))[:16]+'</div></div>',unsafe_allow_html=True)
        except: continue

def render_events():
    for date,name,scope,bg,clr,desc in EVENTS:
        st.markdown('<div class="ev"><div class="ev-date">'+date+'</div>'
            '<div><span class="ev-tag" style="background:'+bg+';color:'+clr+'">'+scope+'</span> '
            '<span style="font-size:12px;font-weight:600;color:#E6EDF3;margin-left:4px">'+name+'</span>'
            '<div class="ev-body">'+desc+'</div></div></div>',unsafe_allow_html=True)

# ═══ MPL PNG (PPTX) ═══
def _h2r(h): h=h.lstrip("#");return tuple(int(h[i:i+2],16)/255 for i in (0,2,4))
def mpl_monthly(m,metal,color):
    fig,ax=plt.subplots(figsize=(9.2,4.2),facecolor=CARD);ax.set_facecolor(CARD)
    l=list(m["Label"]);p=list(m["Price"]);pv=[None]+p[:-1]
    for i,(pr,pp) in enumerate(zip(p,pv)):
        c=(0,.78,.33,.12) if (pp is None or pr>=pp) else (1,.09,.27,.12);ax.bar(i,pr,color=c,width=.75,zorder=1)
    ax.plot(range(len(l)),p,color=_h2r(color),linewidth=2.2,zorder=3)
    for i,(pr,pp) in enumerate(zip(p,pv)):
        ax.scatter(i,pr,color=GR if (pp is None or pr>=pp) else RD,s=28,zorder=4,linewidths=.8,edgecolors=CARD)
        ax.text(i,pr*1.002,f"{pr:,.0f}",ha="center",va="bottom",fontsize=7,color=SC)
    ax.set_xticks(range(len(l)));ax.set_xticklabels(l,rotation=35,ha="right",fontsize=8,color=SC)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:,.0f}"));ax.tick_params(colors=SC,labelsize=8)
    ax.spines[:].set_color(GD);ax.grid(axis="y",color=GD,linewidth=.5,zorder=0);ax.grid(axis="x",visible=False)
    ax.text(0,-.18,"Source: LME / FRED",transform=ax.transAxes,fontsize=7.5,color=MT)
    fig.tight_layout(pad=.8);buf=io.BytesIO();fig.savefig(buf,format="png",dpi=150,facecolor=CARD,bbox_inches="tight");plt.close(fig);buf.seek(0);return buf.read()

def mpl_quarterly(q,metal,color):
    fig,ax=plt.subplots(figsize=(9.2,4.2),facecolor=CARD);ax.set_facecolor(CARD)
    l=list(q["Label"]);p=list(q["Price"]);pv=[None]+p[:-1]
    bc=[GR if (pp is None or pr>=pp) else RD for pr,pp in zip(p,pv)]
    ax.bar(range(len(l)),p,color=bc,alpha=.85,width=.55,zorder=2)
    for i,pr in enumerate(p): ax.text(i,pr*1.003,f"{pr:,.0f}",ha="center",va="bottom",fontsize=8.5,color=TX)
    ax.set_xticks(range(len(l)));ax.set_xticklabels(l,rotation=20,ha="right",fontsize=8,color=SC)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:,.0f}"));ax.tick_params(colors=SC,labelsize=8)
    ax.spines[:].set_color(GD);ax.grid(axis="y",color=GD,linewidth=.5,zorder=0);ax.grid(axis="x",visible=False)
    fig.tight_layout(pad=.8);buf=io.BytesIO();fig.savefig(buf,format="png",dpi=150,facecolor=CARD,bbox_inches="tight");plt.close(fig);buf.seek(0);return buf.read()

# ═══ PPTX ═══
def _rgb(h):h=h.lstrip("#");return RGBColor(int(h[:2],16),int(h[2:4],16),int(h[4:],16))
def _ts(sl,l,t,w,h,text,pt,bold=False,color=TX,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=align
    r=p.add_run();r.text=text;r.font.size=Pt(pt);r.font.bold=bold;r.font.color.rgb=_rgb(color)
def _kp(sl,l,t,w,h,label,value,sub,bg):
    rect=sl.shapes.add_shape(1,l,t,w,h);rect.fill.solid();rect.fill.fore_color.rgb=_rgb(bg);rect.line.color.rgb=_rgb(BD);rect.line.width=Pt(.5)
    _ts(sl,l+Inches(.08),t+Inches(.06),w-Inches(.16),Inches(.2),label,7,color=SC,align=PP_ALIGN.CENTER)
    _ts(sl,l+Inches(.06),t+Inches(.24),w-Inches(.12),Inches(.42),value,15,bold=True,align=PP_ALIGN.CENTER)
    if sub:_ts(sl,l+Inches(.06),t+Inches(.64),w-Inches(.12),Inches(.2),sub,7,color=GR,align=PP_ALIGN.CENTER)
def _sl(prs,title,sub,png_b,kpi,accent):
    sl=prs.slides.add_slide(prs.slide_layouts[6]);W,H=prs.slide_width,prs.slide_height
    bg=sl.shapes.add_shape(1,0,0,W,H);bg.fill.solid();bg.fill.fore_color.rgb=_rgb(BG);bg.line.fill.background()
    hd=sl.shapes.add_shape(1,0,0,W,Inches(.7));hd.fill.solid();hd.fill.fore_color.rgb=_rgb(CARD);hd.line.fill.background()
    ac=sl.shapes.add_shape(1,0,Inches(.7),W,Inches(.022));ac.fill.solid();ac.fill.fore_color.rgb=_rgb(accent);ac.line.fill.background()
    _ts(sl,Inches(.25),Inches(.06),Inches(7),Inches(.36),title,16,bold=True)
    _ts(sl,Inches(.25),Inches(.41),Inches(7),Inches(.22),sub,8.5,color=SC)
    sl.shapes.add_picture(io.BytesIO(png_b),Inches(.12),Inches(.76),width=Inches(6.75),height=Inches(4.15))
    kx,kw,kh=Inches(7.1),Inches(2.65),Inches(1.08)
    _kp(sl,kx,Inches(.78),kw,kh,kpi["l1"],kpi["v1"],None,CARD)
    _kp(sl,kx,Inches(1.96),kw,kh,kpi["l2"],kpi["v2"],None,CARD)
    _kp(sl,kx,Inches(3.14),kw,kh,kpi["l3"],kpi["v3"],kpi.get("s3",""),CARD)
def _kd(m):
    c=m["Price"].iloc[-1];cl=m["Label"].iloc[-1];lo=m["Price"].min();ll=m.loc[m["Price"].idxmin(),"Label"]
    f=m["Price"].iloc[0];fl=m["Label"].iloc[0];ch=c-f;pc=ch/f*100 if f else 0;a="▲" if ch>=0 else "▼"
    return {"l1":f"CURRENT ({cl})","v1":f"${c:,.0f}","l2":f"LOW ({ll})","v2":f"${lo:,.0f}","l3":f"RETURN {fl}→{cl}","v3":f"{a} ${abs(ch):,.0f}","s3":f"{a} {abs(pc):.1f}%"}
def build_pptx(mdata,qdata,sel):
    def f(df):r=df[df["Label"].isin(sel)];return r if not r.empty else df.tail(10)
    def fq(df,mf):
        if mf.empty:return df.tail(4)
        r=df[(df["Date"]>=mf["Date"].min())&(df["Date"]<=mf["Date"].max())];return r if not r.empty else df.tail(4)
    prs=Presentation();prs.slide_width=Inches(10);prs.slide_height=Inches(5.625)
    keys=[k for k in ["cu","al","br","ag"] if k in mdata]
    for key in keys:
        mf=f(mdata[key]);qf=fq(qdata[key],mf)
        _sl(prs,NAMES[key]+" — Monthly Average",UNITS[key]+" · "+mf["Label"].iloc[0]+" → "+mf["Label"].iloc[-1],
            mpl_monthly(mf,NAMES[key],CLR[key]),_kd(mf),CLR[key])
    for key in keys:
        mf=f(mdata[key]);qf=fq(qdata[key],mf);kp=_kd(mf)
        if len(qf):kp["l1"]=f"QTR ({qf['Label'].iloc[-1]})";kp["v1"]=f"${qf['Price'].iloc[-1]:,.0f}"
        _sl(prs,NAMES[key]+" — Quarterly Average",
            (qf["Label"].iloc[0] if len(qf) else "")+" → "+(qf["Label"].iloc[-1] if len(qf) else ""),
            mpl_quarterly(qf,NAMES[key],CLR[key]),kp,CLR[key])
    buf=io.BytesIO();prs.save(buf);buf.seek(0);return buf.read()

# ═══ MAIN ═══
def main():
    with st.spinner(""):
        raw_all={k:load(v) for k,v in SHEETS.items()}

    for k in ["cu","al"]:
        if raw_all[k].empty:
            st.error("No data for "+NAMES[k]+" — check Google Sheets connection."); return

    # Only keep commodities that have REAL data (no synthetic fallback)
    raw_full = {k: v for k, v in raw_all.items() if not v.empty}
    active = list(raw_full.keys())
    missing = [k for k in SHEETS if k not in active]

    global_min = min(raw_full[k]["Date"].min() for k in active)
    global_max = max(raw_full[k]["Date"].max() for k in active)

    with st.sidebar:
        st.markdown('<div style="padding:10px 0"><div style="font-size:24px;font-weight:900;color:#FFF;letter-spacing:3px;font-family:Arial Black,sans-serif;line-height:1">'
            'LME<span style="color:#B87333">.</span></div>'
            '<div style="font-size:9px;color:#8B949E;letter-spacing:2px;text-transform:uppercase;margin-top:3px">Commodities Terminal</div></div>'
            '<div style="height:2px;background:linear-gradient(90deg,#B87333,transparent);margin-bottom:14px"></div>',unsafe_allow_html=True)

        st.markdown('<div style="font-size:10px;color:#8B949E;letter-spacing:1.5px;text-transform:uppercase;margin-bottom:6px">Date Range</div>',unsafe_allow_html=True)
        preset = st.radio("Range", ["All","YTD","6M","3M","1M","Custom"], horizontal=True, label_visibility="collapsed")
        today = global_max
        if preset == "All":   start_d, end_d = global_min.date(), global_max.date()
        elif preset == "YTD": start_d, end_d = datetime(today.year,1,1).date(), today.date()
        elif preset == "6M":  start_d, end_d = (today - pd.DateOffset(months=6)).date(), today.date()
        elif preset == "3M":  start_d, end_d = (today - pd.DateOffset(months=3)).date(), today.date()
        elif preset == "1M":  start_d, end_d = (today - pd.DateOffset(months=1)).date(), today.date()
        else:
            c1c,c2c = st.columns(2)
            with c1c: start_d = st.date_input("From", value=global_min.date(), min_value=global_min.date(), max_value=global_max.date(), label_visibility="collapsed")
            with c2c: end_d   = st.date_input("To", value=global_max.date(), min_value=global_min.date(), max_value=global_max.date(), label_visibility="collapsed")
        st.caption(f"{start_d:%d %b %Y} → {end_d:%d %b %Y}")
        st.markdown('<div style="height:1px;background:#30363D;margin:12px 0"></div>',unsafe_allow_html=True)

    sd = pd.Timestamp(start_d); ed = pd.Timestamp(end_d)
    raw = {k: v[(v["Date"]>=sd)&(v["Date"]<=ed)].reset_index(drop=True) for k,v in raw_full.items()}
    raw = {k: (v if len(v)>=2 else raw_full[k]) for k,v in raw.items()}

    M={k:mavg(raw[k]) for k in active}; Q={k:qavg(raw[k]) for k in active}

    data={k+"_m":M[k] for k in active}
    ticker_bar(data)

    if missing:
        names = ", ".join(NAMES[k] for k in missing)
        st.warning("No data yet for: "+names+". These Google Sheets are empty — populate them to see live charts. (No fake data is shown.)")

    with st.sidebar:
        all_m=sorted(set(sum([M[k]["Label"].tolist() for k in M],[])),key=lambda x:datetime.strptime(x,"%b-%y"))
        sel=st.multiselect("Export Months",all_m,default=all_m[-10:] if len(all_m)>=10 else all_m)
        st.caption(f"{len(active)*2} slides: {len(active)} monthly + {len(active)} quarterly")
        if st.button("⬇ Generate PPTX",type="primary"):
            if not sel:st.warning("Select months.")
            else:
                with st.spinner("Rendering 8 slides…"):
                    try:
                        d=build_pptx(M,Q,sel)
                        st.download_button("⬇ Download",data=d,file_name=f"Commodities_{datetime.now():%Y%m%d}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",width="stretch")
                    except Exception as e:st.error(str(e))
        inr=get_usd_inr()
        st.markdown('<div style="border-top:1px solid #30363D;margin-top:14px;padding-top:10px;font-size:10px;color:#484F58">'
            '<b style="color:#8B949E">USD/INR</b><br><span style="font-size:18px;color:#FFD600;font-weight:700">₹'+f'{inr:.2f}'+'</span></div>',unsafe_allow_html=True)
        short={"cu":"Cu","al":"Al","br":"Brent","ag":"Ag"}
        st.markdown('<div style="font-size:10px;color:#484F58;margin-top:8px"><b style="color:#8B949E">DATA</b><br>'
            +'<br>'.join([short[k]+': '+raw[k]["Date"].min().strftime("%d %b %y")+'→'+raw[k]["Date"].max().strftime("%d %b %y") for k in active])
            +'</div>',unsafe_allow_html=True)

    tab_labels={"cu":"🟤 COPPER","al":"⚙️ ALUMINIUM","br":"🛢️ BRENT","ag":"🪙 SILVER"}
    all_tab_names=[tab_labels[k] for k in active]+["🌍 MARKET INTEL","📊 COMPARE","📋 DATA"]
    tabs=st.tabs(all_tab_names)
    n=len(active)
    intel_tab=tabs[n]; compare_tab=tabs[n+1]; data_tab=tabs[n+2]

    for idx,key in enumerate(active):
        with tabs[idx]:
            c1,c2=st.columns([3.2,1])
            with c1:st.plotly_chart(chart_live(raw[key],NAMES[key],CLR[key],UNITS[key]),width="stretch",key=key+"_l")
            with c2:kpis(M[key],UNITS[key])
            # Data-driven insight banner (works for ALL commodities)
            st.markdown('<div class="sh" style="margin-top:6px">📊 WHY IS THE PRICE MOVING? (AUTO-ANALYSIS)</div>',unsafe_allow_html=True)
            render_data_insights(raw[key],UNITS[key])
            st.markdown('<div class="sh" style="margin-top:10px">QUARTERLY & ROLLING</div>',unsafe_allow_html=True)
            r1,r2=st.columns(2)
            with r1:st.plotly_chart(chart_q(Q[key],lfc(Q[key],2),NAMES[key],CLR[key]),width="stretch",key=key+"_q")
            with r2:st.plotly_chart(chart_roll(raw[key],NAMES[key],CLR[key]),width="stretch",key=key+"_r")
            st.markdown('<div class="sh" style="margin-top:14px">MONTHLY TABLE</div>',unsafe_allow_html=True)
            tbl(M[key])

    # MARKET INTEL
    INTEL_MAP={"cu":"Copper","al":"Aluminium","br":"Brent Oil","ag":"Silver"}
    INTEL_REV={v:k for k,v in INTEL_MAP.items()}
    with intel_tab:
        intel_options=[INTEL_MAP[k] for k in active]
        sel_c=st.selectbox("Commodity",intel_options,label_visibility="collapsed")
        sel_key=INTEL_REV[sel_c]

        # Data-driven auto analysis first (works for every commodity)
        st.markdown('<div class="sh">📊 AUTO-ANALYSIS — WHY THE PRICE IS MOVING</div>',unsafe_allow_html=True)
        render_data_insights(raw[sel_key],UNITS[sel_key])

        st.markdown('<div class="sh" style="margin-top:16px">TOP 5 PRODUCING COUNTRIES</div>',unsafe_allow_html=True)
        render_top5(sel_c)

        i1,i2=st.columns(2)
        with i1:
            st.markdown('<div class="sh" style="margin-top:14px">SUPPLY CHAIN INSIGHTS</div>',unsafe_allow_html=True)
            render_insights(sel_c)
        with i2:
            st.markdown('<div class="sh" style="margin-top:14px">LATEST NEWS</div>',unsafe_allow_html=True)
            render_news(sel_c)

        st.markdown('<div class="sh" style="margin-top:16px">KEY EVENTS CALENDAR</div>',unsafe_allow_html=True)
        render_events()

        # VOLATILITY
        st.markdown('<div class="sh" style="margin-top:16px">30-DAY VOLATILITY</div>',unsafe_allow_html=True)
        fig_v=go.Figure()
        for key in active:
            d=raw[key].copy().sort_values("Date"); d["vol"]=d["Price"].pct_change().rolling(30).std()*100
            fig_v.add_trace(go.Scatter(x=d["Date"],y=d["vol"],line=dict(color=CLR[key],width=1.5),name=NAMES[key],mode="lines"))
        fig_v.update_layout(**_CL,title=dict(text="<b>30-Day Rolling Volatility</b> <span style='color:#8B949E;font-size:11px'>Daily % std dev</span>",font=dict(size=14)),height=300)
        st.plotly_chart(fig_v,width="stretch",key="vol")

        # CORRELATION (needs 2+ commodities)
        if len(active)>=2:
            st.markdown('<div class="sh" style="margin-top:12px">CORRELATION MATRIX</div>',unsafe_allow_html=True)
            merged=raw[active[0]][["Date"]].copy()
            for key in active:
                d=raw[key][["Date","Price"]].rename(columns={"Price":NAMES[key]})
                merged=pd.merge(merged,d,on="Date",how="inner")
            if len(merged)>10:
                corr=merged.drop(columns="Date").corr()
                fig_c=go.Figure(data=go.Heatmap(z=corr.values,x=corr.columns,y=corr.columns,
                    colorscale=[[0,"#FF1744"],[0.5,"#161B22"],[1,"#00C853"]],zmin=-1,zmax=1,
                    text=[[f"{v:.2f}" for v in row] for row in corr.values],texttemplate="%{text}",
                    textfont=dict(size=13,color="#E6EDF3")))
                fig_c.update_layout(paper_bgcolor=CARD,plot_bgcolor=CARD,font=dict(color=TX,size=11),
                    margin=dict(l=10,r=10,t=40,b=10),height=320,
                    title=dict(text="<b>Price Correlation</b>",font=dict(size=14,color=TX)))
                st.plotly_chart(fig_c,width="stretch",key="corr")

        # INR LANDED COST
        st.markdown('<div class="sh" style="margin-top:12px">INR LANDED COST</div>',unsafe_allow_html=True)
        inr=get_usd_inr()
        cols=st.columns(len(active))
        for col,key in zip(cols,active):
            cur=M[key]["Price"].iloc[-1]; inr_val=cur*inr
            with col:
                st.markdown(
                    '<div class="pc" style="text-align:center">'
                    '<div style="font-size:10px;color:#8B949E;letter-spacing:1px;text-transform:uppercase">'+NAMES[key]+'</div>'
                    '<div style="font-size:20px;font-weight:700;color:#FFD600;margin:6px 0">₹'+f'{inr_val:,.0f}'+'</div>'
                    '<div style="font-size:10px;color:#484F58">$'+f'{cur:,.2f}'+' × ₹'+f'{inr:.2f}'+'</div>'
                    '</div>',unsafe_allow_html=True)

    # COMPARE
    with compare_tab:
        st.markdown('<div class="sh">INDEXED PERFORMANCE (Base=100)</div>',unsafe_allow_html=True)
        fig=go.Figure()
        for key in active:
            d=M[key].copy();d["I"]=d["Price"]/d["Price"].iloc[0]*100
            fig.add_trace(go.Scatter(x=d["Date"],y=d["I"],line=dict(color=CLR[key],width=2),name=NAMES[key]))
        fig.add_hline(y=100,line_dash="dot",line_color=MT,line_width=1)
        fig.update_layout(**_CL,title=dict(text="<b>Indexed Comparison</b>",font=dict(size=15)),height=380)
        st.plotly_chart(fig,width="stretch",key="comp")

    # DATA
    with data_tab:
        cols=st.columns(len(active))
        for col,key in zip(cols,active):
            with col:
                st.markdown('<div class="sh">'+NAMES[key]+'</div>',unsafe_allow_html=True)
                dp=raw[key][["Date","Price"]].copy();dp["Date"]=dp["Date"].dt.strftime("%d %b %Y");dp["Price"]=dp["Price"].apply(lambda x:f"${x:,.2f}")
                st.dataframe(dp,height=400,width="stretch",hide_index=True)

if __name__=="__main__":
    main()
